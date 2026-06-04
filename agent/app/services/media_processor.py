"""
Media processing service.

Handles: audio transcription, image description, PDF/Word/Excel/PPT extraction.
All methods return extracted text ready to be injected into the conversation.
"""
import io
import httpx
from loguru import logger
from openai import AsyncOpenAI
from app.core.config import get_settings

settings = get_settings()
_openai = AsyncOpenAI(api_key=settings.openai_api_key)


async def _download(url: str) -> bytes:
    async with httpx.AsyncClient(timeout=60) as client:
        resp = await client.get(url)
        resp.raise_for_status()
        return resp.content


async def transcribe_audio(url: str) -> str:
    """Transcribe audio file via OpenAI Whisper."""
    try:
        data = await _download(url)
        transcript = await _openai.audio.transcriptions.create(
            model="whisper-1",
            file=("audio.ogg", io.BytesIO(data), "audio/ogg"),
        )
        return f"[Áudio transcrito]: {transcript.text}"
    except Exception as e:
        logger.error(f"Audio transcription failed: {e}")
        return "[Áudio não processado]"


async def describe_image(url: str) -> str:
    """Describe image via GPT-4o Vision."""
    try:
        response = await _openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": "Descreva esta imagem detalhadamente."},
                    {"type": "image_url", "image_url": {"url": url}},
                ],
            }],
            max_tokens=500,
        )
        description = response.choices[0].message.content or ""
        return f"[Imagem recebida]: {description}"
    except Exception as e:
        logger.error(f"Image description failed: {e}")
        return "[Imagem não processada]"


async def extract_pdf(url: str) -> str:
    """Extract text from PDF."""
    try:
        import pymupdf
        data = await _download(url)
        doc = pymupdf.open(stream=data, filetype="pdf")
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return f"[PDF recebido]:\n{text[:3000]}"
    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        return "[PDF não processado]"


async def extract_docx(url: str) -> str:
    """Extract text from Word document."""
    try:
        from docx import Document
        data = await _download(url)
        doc = Document(io.BytesIO(data))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return f"[Documento Word recebido]:\n{text[:3000]}"
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        return "[Documento não processado]"


async def extract_xlsx(url: str) -> str:
    """Extract text from Excel spreadsheet."""
    try:
        import openpyxl
        data = await _download(url)
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        rows = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                line = " | ".join(str(c) for c in row if c is not None)
                if line:
                    rows.append(line)
        text = "\n".join(rows[:100])
        return f"[Planilha Excel recebida]:\n{text}"
    except Exception as e:
        logger.error(f"XLSX extraction failed: {e}")
        return "[Planilha não processada]"


async def extract_pptx(url: str) -> str:
    """Extract text from PowerPoint."""
    try:
        from pptx import Presentation
        data = await _download(url)
        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [s.text for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
            if texts:
                slides.append(f"Slide {i}: {' | '.join(texts)}")
        return f"[Apresentação recebida]:\n{chr(10).join(slides[:30])}"
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        return "[Apresentação não processada]"


async def process_media(media_url: str, media_type: str, mimetype: str = "") -> str:
    """Route media to the correct processor based on type/mimetype."""
    if media_type == "audio":
        return await transcribe_audio(media_url)
    if media_type == "image":
        return await describe_image(media_url)
    if media_type == "document":
        if "pdf" in mimetype:
            return await extract_pdf(media_url)
        if "word" in mimetype or "docx" in mimetype:
            return await extract_docx(media_url)
        if "excel" in mimetype or "xlsx" in mimetype or "spreadsheet" in mimetype:
            return await extract_xlsx(media_url)
        if "powerpoint" in mimetype or "pptx" in mimetype or "presentation" in mimetype:
            return await extract_pptx(media_url)
        return await extract_pdf(media_url)  # fallback
    return ""
